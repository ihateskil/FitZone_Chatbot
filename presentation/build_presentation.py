"""
FitZone Graduation Project Presentation — 15 slides
Dark theme (#0E0E10), green accent (#39D353), academic committee tone.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ─────────────────────────────────────────────────
BG = RGBColor(0x0E, 0x0E, 0x10)
GREEN = RGBColor(0x39, 0xD3, 0x53)
BLUE = RGBColor(0x37, 0x8A, 0xDD)
AMBER = RGBColor(0xEF, 0x9F, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0xAA, 0xAA, 0xAA)
DARK_CARD = RGBColor(0x1A, 0x1A, 0x1E)
DARK_CARD2 = RGBColor(0x14, 0x14, 0x18)

PRES_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(PRES_DIR, "FitZone_Presentation.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helper functions ───────────────────────────────────────────────
def _bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, fill=DARK_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def _text_box(slide, left, top, width, height, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Inter"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def _add_para(tf, text, size=14, color=WHITE, bold=False, space_before=Pt(6), font="Inter", align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.space_before = space_before
    p.alignment = align
    return p


def _tag(slide, left, top, text, fill=GREEN, text_color=BG, size=10):
    shape = _rect(slide, left, top, Inches(len(text) * 0.18 + 0.3), Inches(0.3), fill=fill, radius=0.05)
    _text_box(slide, left + Inches(0.05), top + Inches(0.02), shape.width - Inches(0.1), Inches(0.26),
              text, size=size, color=text_color, bold=True, align=PP_ALIGN.CENTER)


def _accent_line(slide, left, top, width=Inches(1.2)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()


def _card_grid(slide, cards, start_top, cols=3, card_w=Inches(3.8), card_h=Inches(1.6), gap=Inches(0.3)):
    """cards = list of (title, body, tag_text, tag_color) where tag_color is RGBColor"""
    x_start = Inches(0.6)
    for i, (title, body, tag_t, tag_c) in enumerate(cards):
        col = i % cols
        row = i // cols
        x = x_start + col * (card_w + gap)
        y = start_top + row * (card_h + gap)
        card = _rect(slide, x, y, card_w, card_h, fill=DARK_CARD)
        _tag(slide, x + Inches(0.15), y + Inches(0.12), tag_t, fill=tag_c)
        _text_box(slide, x + Inches(0.15), y + Inches(0.45), card_w - Inches(0.3), Inches(0.35),
                  title, size=14, color=WHITE, bold=True)
        _text_box(slide, x + Inches(0.15), y + Inches(0.85), card_w - Inches(0.3), Inches(0.65),
                  body, size=10, color=GREY)


# ── Build slides ───────────────────────────────────────────────────

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_bg(sl)

# Green accent bar top
bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Pt(6))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()

# FitZone wordmark
_text_box(sl, Inches(1), Inches(1.0), Inches(5), Inches(0.8),
          "FITZONE", size=44, color=GREEN, bold=True)

_sub = _text_box(sl, Inches(1), Inches(1.7), Inches(5), Inches(0.5),
          "Integrated Online Fitness & Coaching Management Platform", size=18, color=GREY)

# University
_text_box(sl, Inches(1), Inches(3.2), Inches(7), Inches(0.4),
          "Faculty of Computers & Artificial Intelligence", size=16, color=WHITE, bold=True)
_text_box(sl, Inches(1), Inches(3.6), Inches(7), Inches(0.4),
          "Assiut National University", size=14, color=GREY)

# Supervisors
_text_box(sl, Inches(1), Inches(4.4), Inches(7), Inches(0.3),
          "Supervisors:", size=12, color=GREY)
_text_box(sl, Inches(1), Inches(4.75), Inches(7), Inches(0.3),
          "Dr. Amal AbdelAzim  ·  Eng. Nehad AbdelRahman", size=14, color=WHITE)

# Date
_text_box(sl, Inches(1), Inches(5.4), Inches(7), Inches(0.3),
          "June 2026", size=12, color=GREY)

# Team members (right side)
tm_x = Inches(7.2)
_t = _text_box(sl, tm_x, Inches(1.2), Inches(5.2), Inches(5.5),
          "", size=12, color=GREY)
_t.text_frame.word_wrap = True
_t.text_frame.paragraphs[0].text = "TEAM MEMBERS"
_t.text_frame.paragraphs[0].font.size = Pt(10); _t.text_frame.paragraphs[0].font.bold = True; _t.text_frame.paragraphs[0].font.color.rgb = GREEN
members = [
    "Mohamed Osamy Mohamed — Backend Architecture & Proposal Engine",
    "Ahmed Sayed Khafaga — Authentication & Security Layer",
    "Esraa Khamis Adwi — Nutrition Plan System (Backend)",
    "Ahmed AbdelFattah Saad — Workout Program Builder (Backend)",
    "Khaled Ayman Alwan — Frontend — Coach Portal",
    "Mohamed Akram AbdelLatif — Frontend — Trainee Portal",
    "Kholoud Ibrahim Ahmed — Frontend — Admin Portal & UI",
    "Khaled Ibrahim AbdelMonem — SignalR Chat & Payments",
]
for m in members:
    _add_para(_t.text_frame, m, size=11, color=WHITE, space_before=Pt(4))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Problem Statement", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Why existing fitness platforms fall short", size=14, color=GREY)

pain_points = [
    ("Fragmented Ecosystem", "No single platform connects trainees, coaches, and administrators — users juggle 3-4 separate tools for workouts, nutrition, payments, and communication.", GREEN),
    ("Manual Nutrition Tracking", "Nutrition logging is tedious, error-prone, and lacks intelligent synchronization with training demands and metabolic calculations.", BLUE),
    ("No Unified Portal", "Coaches manage clients across spreadsheets, messaging apps, and disconnected billing systems — no centralized dashboard exists.", AMBER),
    ("Disconnected Payments", "Subscription billing, coach commissions, and session fees operate in silos with no integration to the coaching workflow itself.", RGBColor(0xE0, 0x68, 0x59)),
    ("No Feedback Loop", "Trainees receive static plans with no mechanism for coaches to assess progress, adjust macros, or unlock the next training week intelligently.", RGBColor(0x7B, 0x68, 0xEE)),
]

for i, (title, desc, color) in enumerate(pain_points):
    y = Inches(1.9) + i * Inches(1.05)
    # Number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.45), Inches(0.45))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    _text_box(sl, Inches(0.6), y + Inches(0.06), Inches(0.45), Inches(0.35),
              str(i + 1), size=14, color=BG, bold=True, align=PP_ALIGN.CENTER)
    _text_box(sl, Inches(1.2), y + Inches(0.02), Inches(3.5), Inches(0.3),
              title, size=14, color=WHITE, bold=True)
    _text_box(sl, Inches(1.2), y + Inches(0.35), Inches(10.5), Inches(0.6),
              desc, size=10, color=GREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Our Solution — Three Role-Based Portals", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "One platform, three perspectives — unified by a shared data layer and business logic.", size=14, color=GREY)

portals = [
    ("Admin Portal", "System-wide oversight, platform analytics, coach & trainee management, payment reconciliation, content libraries, and audit logging.", GREEN, "ADMIN"),
    ("Coach Portal", "Client roster, workout program builder, nutrition plan designer, check-in review, proposal approval, real-time chat, and commission tracking.", BLUE, "COACH"),
    ("Trainee Portal", "Personalized workout & nutrition plans, daily check-ins, progress tracking, direct messaging, subscription management, and week-unlock flow.", AMBER, "TRAINEE"),
]
for i, (title, desc, color, tag) in enumerate(portals):
    x = Inches(0.6) + i * Inches(4.2)
    card = _rect(sl, x, Inches(1.8), Inches(3.9), Inches(4.5), fill=DARK_CARD)
    _tag(sl, x + Inches(0.2), Inches(2.0), tag, fill=color)
    _text_box(sl, x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(0.4),
              title, size=20, color=WHITE, bold=True)
    _text_box(sl, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(2.0),
              desc, size=12, color=GREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Market Analysis
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Market Analysis", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Competitor landscape and our unique value proposition", size=14, color=GREY)

# Comparison table
cols_x = [Inches(0.6), Inches(3.0), Inches(5.0), Inches(7.0), Inches(9.0), Inches(11.0)]
col_w = [Inches(2.2), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)]
headers = ["", "MyFitnessPal", "Trainerize", "TrueCoach", "FitZone"]

# Header row
for j, h in enumerate(headers):
    c = GREEN if h == "FitZone" else WHITE
    _text_box(sl, cols_x[j], Inches(1.7), col_w[j], Inches(0.35),
              h, size=12, color=c, bold=True, align=PP_ALIGN.CENTER)

# Features + checkmarks
features = [
    ("Role-Based Portals", "✗", "✗", "✗", "✓"),
    ("AI Proposal Engine", "✗", "✗", "✗", "✓"),
    ("Nutrition System", "✓", "✗", "✗", "✓"),
    ("Built-in Payments", "✗", "✓", "✓", "✓"),
    ("Real-Time Chat", "✗", "✗", "✗", "✓"),
    ("Week-Unlock Flow", "✗", "✗", "✗", "✓"),
    ("Unified Dashboard", "✗", "✓", "✓", "✓"),
]

for i, row in enumerate(features):
    y = Inches(2.2) + i * Inches(0.45)
    for j, val in enumerate(row):
        color = WHITE if j == 0 else (GREEN if val == "✓" else AMBER)
        _text_box(sl, cols_x[j], y, col_w[j], Inches(0.35),
                  val, size=11, color=color, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# UVP box
uvp = _rect(sl, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.2), fill=DARK_CARD)
_uvp_text = _text_box(sl, Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.9),
          "", size=12, color=WHITE)
_uvp_text.text_frame.word_wrap = True
_uvp_text.text_frame.paragraphs[0].text = "Unique Value Proposition"
_uvp_text.text_frame.paragraphs[0].font.bold = True
_uvp_text.text_frame.paragraphs[0].font.color.rgb = GREEN
_uvp_text.text_frame.paragraphs[0].font.size = Pt(14)
_add_para(_uvp_text.text_frame,
  "FitZone is the only platform combining role-based portals, an AI-driven proposal engine, and an integrated payment ecosystem "
  "in a single platform. Where competitors offer point solutions, FitZone delivers an end-to-end coaching lifecycle — "
  "from enrollment and assessment through week-unlock progression and performance analytics.",
  size=11, color=GREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Target Audience
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Target Audience — Personas", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Three user archetypes driving every design decision", size=14, color=GREY)

personas = [
    ("Sara — Platform Admin", "30, Gym Chain Operations Manager", GREEN, [
        "Manages 12 coaches and 200+ trainees across 3 locations",
        "Needs payment reconciliation, content library control, user analytics",
        "Pain: Spreadsheets for everything — no unified oversight",
    ]),
    ("Ahmed — Fitness Coach", "35, Certified Personal Trainer", BLUE, [
        "Works with 25 active clients, creates custom programs weekly",
        "Needs workout builder, nutrition planner, check-in reviews, commission tracking",
        "Pain: Juggles WhatsApp, Excel, and separate billing tools",
    ]),
    ("Fatima — Trainee", "24, Software Engineer, fitness enthusiast", AMBER, [
        "Follows a structured training program with macro targets",
        "Needs daily workouts, meal guidance, progress photos, direct coach chat",
        "Pain: Previous apps had generic plans — no personalization or coach feedback",
    ]),
]

for i, (name, subtitle, color, bullets) in enumerate(personas):
    x = Inches(0.6) + i * Inches(4.2)
    card = _rect(sl, x, Inches(1.8), Inches(3.9), Inches(5.0), fill=DARK_CARD)
    _tag(sl, x + Inches(0.2), Inches(2.0), "PERSONA", fill=color)
    _text_box(sl, x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(0.4),
              name, size=16, color=WHITE, bold=True)
    _text_box(sl, x + Inches(0.2), Inches(2.9), Inches(3.5), Inches(0.3),
              subtitle, size=10, color=GREY)
    for j, b in enumerate(bullets):
        _text_box(sl, x + Inches(0.2), Inches(3.5 + j * 0.5), Inches(3.5), Inches(0.5),
                  f"• {b}", size=10, color=WHITE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Core Features
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Core Features", size=28, color=WHITE, bold=True)

features = [
    ("JWT Authentication", "Secure login with role-based access control across all three portals.", "AUTH", GREEN),
    ("Stripe Payments", "Subscription billing, coach commissions, and payout management integrated end-to-end.", "PAY", GREEN),
    ("Workout Builder", "Drag-and-drop program construction with exercise libraries, sets/reps/RPE configuration.", "WORKOUT", BLUE),
    ("Nutrition System", "Automated macro calculation from TDEE, coach-approved proposals, weekly propagation.", "NUTRITION", BLUE),
    ("AI Proposal Engine", "Check-in analysis → TDEE calc → macro proposal → coach approval → auto-apply.", "AI", AMBER),
    ("SignalR Chat", "Real-time messaging between trainees and coaches with read receipts and notifications.", "CHAT", AMBER),
    ("Week Unlock", "Coaches review check-in data before granting access to the next training week.", "WEEK", RGBColor(0x7B,0x68,0xEE)),
    ("Content Libraries", "Admin-managed exercise, meal, and formula reference libraries used across all programs.", "LIBS", RGBColor(0x7B,0x68,0xEE)),
    ("Progress Tracking", "Weight logs, progress photos, performance metrics, and visual trend charts per trainee.", "TRACK", RGBColor(0xE0,0x68,0x59)),
]
_card_grid(sl, features, Inches(1.4), cols=3, card_w=Inches(3.8), card_h=Inches(1.7))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — AI Proposal Engine
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "AI Proposal Engine", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Automated macro proposal workflow — from check-in to coach approval", size=14, color=GREY)

# Flow chart (simplified with boxes and arrows)
steps = [
    "Trainee\nCheck-In", "TDEE\nCalculation", "AI\nProposal", "Coach\nReview", "Auto-Apply\nMacros"
]
colors = [AMBER, GREEN, BLUE, AMBER, GREEN]
for i, (step, c) in enumerate(zip(steps, colors)):
    x = Inches(0.8) + i * Inches(2.5)
    box = _rect(sl, x, Inches(2.0), Inches(2.0), Inches(2.0), fill=DARK_CARD)
    _text_box(sl, x + Inches(0.1), Inches(2.5), Inches(1.8), Inches(1.0),
              step, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(steps) - 1:
        _text_box(sl, x + Inches(1.8), Inches(2.8), Inches(0.7), Inches(0.3),
                  "→", size=20, color=c, bold=True, align=PP_ALIGN.CENTER)

# Description below
desc = _rect(sl, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), fill=DARK_CARD)
_desc_t = _text_box(sl, Inches(1.2), Inches(4.7), Inches(11.0), Inches(2.1),
          "", size=12, color=WHITE)
_desc_t.text_frame.word_wrap = True
_desc_t.text_frame.paragraphs[0].text = "How It Works"
_desc_t.text_frame.paragraphs[0].font.bold = True
_desc_t.text_frame.paragraphs[0].font.color.rgb = GREEN
_desc_t.text_frame.paragraphs[0].font.size = Pt(14)
details = [
    "1. Trainee submits a daily check-in with weight, energy level, sleep, and adherence rating.",
    "2. The engine calculates TDEE using the Mifflin-St Jeor formula × activity multiplier from current week volume.",
    "3. The AI proposes new macros (protein, carbs, fat) ± calorie surplus/deficit aligned to the trainee's goal (cut/bulk/recomp).",
    "4. The coach receives the proposal in their dashboard — one tap to approve, adjust, or reject with notes.",
    "5. Approved macros propagate automatically to the trainee's nutrition plan for the upcoming week.",
]
for d in details:
    _add_para(_desc_t.text_frame, d, size=11, color=GREY, space_before=Pt(4))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — User Journey
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "User Journey — The Coaching Lifecycle", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "From registration to weekly progression — a continuous feedback loop", size=14, color=GREY)

journey = [
    ("1", "Register", "Trainee signs up,\nselects a coach,\nchooses a plan", GREEN),
    ("2", "Enroll", "Coach assigns\ninitial programs\n& nutrition targets", BLUE),
    ("3", "Week 1", "Trainee follows\nworkouts, logs\ncheck-ins daily", AMBER),
    ("4", "Check-In", "Weight, energy,\nadherence data\nsent to coach", RGBColor(0xE0,0x68,0x59)),
    ("5", "AI Proposes", "Engine generates\nmacro adjustments\nfor coach review", RGBColor(0x7B,0x68,0xEE)),
    ("6", "Next Week", "Coach approves →\nweek unlocks →\ncycle repeats", GREEN),
]

for i, (num, title, desc, color) in enumerate(journey):
    angle = i * 60 - 90
    import math
    cx, cy = Inches(6.5), Inches(3.8)
    radius = Inches(2.8)
    rad = math.radians(angle)
    x = cx + radius * math.cos(rad) - Inches(1.0)
    y = cy + radius * math.sin(rad) - Inches(0.7)

    card = _rect(sl, x, y, Inches(2.0), Inches(1.4), fill=DARK_CARD)
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), y - Inches(0.2), Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    _text_box(sl, x + Inches(0.75), y - Inches(0.14), Inches(0.5), Inches(0.4),
              num, size=16, color=BG, bold=True, align=PP_ALIGN.CENTER)
    _text_box(sl, x + Inches(0.1), y + Inches(0.35), Inches(1.8), Inches(0.3),
              title, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _text_box(sl, x + Inches(0.1), y + Inches(0.65), Inches(1.8), Inches(0.7),
              desc, size=9, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Technical Architecture
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Technical Architecture", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Three-tier architecture — modern, scalable, and secure", size=14, color=GREY)

tiers = [
    ("PRESENTATION LAYER", "React.js · TypeScript · Tailwind CSS · shadcn/ui", GREEN, [
        "Single-page application with role-based routing",
        "Responsive design — desktop-first with mobile adaptation",
        "Real-time chat via SignalR JavaScript client",
        "JWT stored in HttpOnly cookies for XSS protection",
    ]),
    ("APPLICATION LAYER", "ASP.NET Core 8 · C# · Entity Framework Core", BLUE, [
        "RESTful API with JWT authentication and role-based authorization",
        "SignalR hubs for real-time messaging and notifications",
        "Stripe API integration for subscription and commission payments",
        "AI proposal engine — server-side TDEE and macro computation",
    ]),
    ("DATA LAYER", "SQL Server · Azure SQL · EF Core Migrations", AMBER, [
        "Normalized relational schema with 20+ entities",
        "Parameterized queries preventing SQL injection",
        "Health data encrypted at rest and isolated per tenant",
        "Migrations for schema versioning and rollback",
    ]),
]

for i, (title, subtitle, color, bullets) in enumerate(tiers):
    x = Inches(0.6) + i * Inches(4.2)
    card = _rect(sl, x, Inches(1.8), Inches(3.9), Inches(5.0), fill=DARK_CARD)
    _tag(sl, x + Inches(0.2), Inches(2.0), title, fill=color)
    _text_box(sl, x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(0.4),
              subtitle, size=10, color=WHITE, bold=True)
    for j, b in enumerate(bullets):
        _text_box(sl, x + Inches(0.2), Inches(3.2 + j * 0.55), Inches(3.5), Inches(0.5),
                  f"• {b}", size=10, color=GREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Security & Privacy
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Security & Privacy", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Protecting user data at every layer of the stack", size=14, color=GREY)

items = [
    ("JWT Authentication", "Stateless tokens with configurable expiry. Role claims (Admin, Coach, Trainee) enforced at every endpoint via custom authorization policies.", GREEN),
    ("Role-Based Access Control", "Three distinct RBAC roles mapped to portal routes, API endpoints, and data visibility. Coaches see only their clients; trainees see only their own data.", BLUE),
    ("HTTPS Everywhere", "All communication encrypted via TLS 1.3. HSTS headers prevent protocol downgrade attacks in production.", RGBColor(0x7B,0x68,0xEE)),
    ("Stripe PCI Compliance", "All payment processing delegated to Stripe Elements. No raw card data touches our servers — PCI Scope is reduced to SAQ A.", AMBER),
    ("Parameterized Queries", "Entity Framework Core generates parameterized SQL — eliminating SQL injection risk across all data access layers.", RGBColor(0xE0,0x68,0x59)),
    ("Health Data Isolation", "Personally identifiable information and health metrics encrypted at rest using AES-256. Trainees can request data export or deletion.", GREEN),
]

for i, (title, desc, color) in enumerate(items):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + col * Inches(4.2)
    y = Inches(1.8) + row * Inches(2.7)
    card = _rect(sl, x, y, Inches(3.9), Inches(2.4), fill=DARK_CARD)
    _tag(sl, x + Inches(0.2), y + Inches(0.12), "SECURITY", fill=color)
    _text_box(sl, x + Inches(0.2), y + Inches(0.55), Inches(3.5), Inches(0.35),
              title, size=14, color=WHITE, bold=True)
    _text_box(sl, x + Inches(0.2), y + Inches(1.0), Inches(3.5), Inches(1.2),
              desc, size=10, color=GREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Testing & Validation
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Testing & Validation", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Comprehensive validation across all system layers", size=14, color=GREY)

tests = [
    ("API Testing", "Postman collection with 80+ test cases covering all endpoints. Automated newman runs in CI pipeline.", GREEN, "POSTMAN"),
    ("Business Logic", "xUnit tests for TDEE calculation, macro propagation, week-unlock rules, and commission arithmetic.", BLUE, "xUNIT"),
    ("Payment Sandbox", "Stripe test mode with 50+ transaction scenarios — subscriptions, prorations, refunds, commission payouts.", AMBER, "STRIPE"),
    ("SignalR Testing", "Multi-client connection tests using SignalR's in-memory backplane. Latency under 200ms at 50 concurrent connections.", RGBColor(0x7B,0x68,0xEE), "SIGNALR"),
    ("End-to-End Simulation", "Full 4-week coaching cycle simulated with synthetic data — registration → enrollment → check-ins → proposal → unlock.", GREEN, "E2E"),
]

for i, (title, desc, color, tag) in enumerate(tests):
    x = Inches(0.6)
    y = Inches(1.8) + i * Inches(1.05)
    card = _rect(sl, x, y, Inches(12.1), Inches(0.9), fill=DARK_CARD)
    _tag(sl, x + Inches(0.2), y + Inches(0.12), tag, fill=color)
    _text_box(sl, x + Inches(1.6), y + Inches(0.1), Inches(3.5), Inches(0.3),
              title, size=13, color=WHITE, bold=True)
    _text_box(sl, x + Inches(1.6), y + Inches(0.4), Inches(10.0), Inches(0.4),
              desc, size=10, color=GREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Challenges & Solutions
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Challenges & Solutions", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Key technical hurdles and how we overcame them", size=14, color=GREY)

challenges = [
    ("Dual-Gate Week Unlock", AMBER,
     "Challenge: A trainee's next week must unlock only when both subscription is active AND coach has approved — but payments and program logic live in separate domains.",
     "Solution: Implemented a composite unlock gate using a service-layer orchestrator that validates subscription status via Stripe API, then checks coach approval in the local database before returning a single boolean."),
    ("TDEE Algorithm Precision", BLUE,
     "Challenge: Mifflin-St Jeor BMR alone is insufficient — activity multipliers vary weekly based on actual training volume, not static self-reporting.",
     "Solution: Introduced a dynamic activity factor derived from the trainee's logged workout volume (sets × reps × weight) averaged over the current week, mapped to a precision scale of 1.2–1.9."),
    ("Macro Propagation Logic", GREEN,
     "Challenge: When a coach approves new macros, changes must cascade across all days of the upcoming week without overwriting manual per-day overrides.",
     "Solution: Implemented a three-way merge strategy: coach-approved baseline → day-level overrides (stored as diffs) → final plan. Overrides survive re-propagation unless explicitly reset."),
    ("Team Coordination", RGBColor(0x7B,0x68,0xEE),
     "Challenge: 8 team members across backend, frontend, and integration work — API contracts, shared types, and data shapes frequently drifted.",
     "Solution: Defined a shared TypeScript types package published as a git submodule. All C# DTOs and TypeScript interfaces generated from a single OpenAPI spec. Weekly integration syncs caught drift early."),
]

for i, (title, color, challenge, solution) in enumerate(challenges):
    y = Inches(1.8) + i * Inches(1.35)
    card = _rect(sl, Inches(0.6), y, Inches(12.1), Inches(1.2), fill=DARK_CARD)
    # Challenge side
    _tag(sl, Inches(0.8), y + Inches(0.08), "CHALLENGE", fill=color)
    _text_box(sl, Inches(0.8), y + Inches(0.4), Inches(11.5), Inches(0.5),
              challenge, size=9, color=WHITE)
    # Solution side
    _tag(sl, Inches(0.8), y + Inches(0.65), "SOLUTION", fill=GREEN)
    _text_box(sl, Inches(0.8), y + Inches(0.95), Inches(11.5), Inches(0.3),
              solution, size=9, color=GREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Monetization
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Monetization Model", size=28, color=WHITE, bold=True)
_text_box(sl, Inches(0.6), Inches(1.2), Inches(10), Inches(0.4),
          "Three revenue streams powering sustainable growth", size=14, color=GREY)

models = [
    ("Membership\nSubscriptions", GREEN, [
        "Monthly / quarterly / annual plans for trainees",
        "Tiered pricing: Basic (workouts only), Pro (workouts + nutrition), Elite (full AI + 1:1 chat)",
        "Recurring billing via Stripe with automated proration and dunning",
    ]),
    ("Coach\nCommissions", BLUE, [
        "Platform takes a percentage (15–25%) of coach service fees",
        "Coaches set their own rates; platform handles invoicing and payouts",
        "Weekly commission reports with Stripe Connect payout automation",
    ]),
    ("Future:\nE-Commerce", AMBER, [
        "In-app supplement store with affiliate partnerships",
        "Premium workout programs & nutrition templates available for purchase",
        "Branded merchandise and equipment partnerships",
    ]),
]

for i, (title, color, bullets) in enumerate(models):
    x = Inches(0.6) + i * Inches(4.2)
    card = _rect(sl, x, Inches(1.8), Inches(3.9), Inches(4.8), fill=DARK_CARD)
    _text_box(sl, x + Inches(0.2), Inches(2.0), Inches(3.5), Inches(0.8),
              title, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        _text_box(sl, x + Inches(0.2), Inches(3.0 + j * 0.6), Inches(3.5), Inches(0.6),
                  f"• {b}", size=10, color=WHITE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — Future Roadmap
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)
_accent_line(sl, Inches(0.6), Inches(0.5))
_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Future Roadmap", size=28, color=WHITE, bold=True)

phases = [
    ("PHASE 1\n✓ COMPLETED", GREEN, [
        "Core platform with 3 role-based portals",
        "JWT authentication & RBAC authorization",
        "Stripe subscription & commission payments",
        "Workout builder & nutrition plan system",
        "AI proposal engine for macro recommendations",
        "SignalR real-time chat",
        "Week-unlock progression flow",
        "Admin content libraries & analytics",
    ]),
    ("PHASE 2\nIN PROGRESS", BLUE, [
        "AI chatbot integration for instant trainee support",
        "E-commerce module for supplements & programs",
        "Advanced analytics dashboard with coach insights",
        "Group training & class scheduling",
        "Multi-language support (Arabic / English)",
    ]),
    ("PHASE 3\nFUTURE", AMBER, [
        "Mobile application (React Native)",
        "Gym equipment & wearable API integration",
        "Barcode scanning for nutrition logging",
        "Video exercise demonstration library",
        "Social features — challenges, leaderboards",
    ]),
]

for i, (title, color, items) in enumerate(phases):
    x = Inches(0.6) + i * Inches(4.2)
    card = _rect(sl, x, Inches(1.5), Inches(3.9), Inches(5.5), fill=DARK_CARD)
    _text_box(sl, x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.7),
              title, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        _text_box(sl, x + Inches(0.2), Inches(2.6 + j * 0.45), Inches(3.5), Inches(0.4),
                  f"• {item}", size=10, color=WHITE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — Team & Closing
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl)

# Green accent bar top
bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Pt(6))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()

_text_box(sl, Inches(0.6), Inches(0.6), Inches(8), Inches(0.6),
          "Team & Closing", size=28, color=WHITE, bold=True)

# Team grid
team = [
    ("Mohamed Osamy", "Backend Architecture\n& Proposal Engine", GREEN),
    ("Ahmed Sayed", "Authentication\n& Security Layer", BLUE),
    ("Esraa Khamis", "Nutrition Plan\nSystem (Backend)", AMBER),
    ("Ahmed AbdelFattah", "Workout Program\nBuilder (Backend)", RGBColor(0xE0,0x68,0x59)),
    ("Khaled Ayman", "Frontend —\nCoach Portal", RGBColor(0x7B,0x68,0xEE)),
    ("Mohamed Akram", "Frontend —\nTrainee Portal", GREEN),
    ("Kholoud Ibrahim", "Frontend — Admin\nPortal & UI System", BLUE),
    ("Khaled Ibrahim", "SignalR Chat\n& Payments", AMBER),
]

cols = 4
for i, (name, role, color) in enumerate(team):
    col = i % cols
    row = i // cols
    x = Inches(0.6) + col * Inches(3.1)
    y = Inches(1.5) + row * Inches(1.5)
    card = _rect(sl, x, y, Inches(2.8), Inches(1.3), fill=DARK_CARD)
    # Dot indicator
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    _text_box(sl, x + Inches(0.35), y + Inches(0.1), Inches(2.3), Inches(0.3),
              name, size=12, color=WHITE, bold=True)
    _text_box(sl, x + Inches(0.35), y + Inches(0.45), Inches(2.3), Inches(0.7),
              role, size=9, color=GREY)

# Closing quote
quote_box = _rect(sl, Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.2), fill=DARK_CARD)
_qt = _text_box(sl, Inches(1.9), Inches(5.0), Inches(9.5), Inches(0.8),
          "", size=18, color=WHITE)
_qt.text_frame.word_wrap = True
_qt.text_frame.paragraphs[0].text = '"Building the future of connected fitness — one macro, one workout, one conversation at a time."'
_qt.text_frame.paragraphs[0].font.italic = True
_qt.text_frame.paragraphs[0].font.color.rgb = GREEN
_qt.text_frame.paragraphs[0].font.size = Pt(20)
_qt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Q&A
_qa = _text_box(sl, Inches(4.5), Inches(6.2), Inches(4.3), Inches(0.5),
          "Thank You — Questions?", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Supervisors at bottom
_text_box(sl, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
          "Supervisors: Dr. Amal AbdelAzim  ·  Eng. Nehad AbdelRahman     |     Faculty of Computers & AI, Assiut National University     |     June 2026",
          size=9, color=GREY, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✅ Saved to {OUTPUT}")
